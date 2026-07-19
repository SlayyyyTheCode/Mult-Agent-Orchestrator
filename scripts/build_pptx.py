#!/usr/bin/env python3
"""Render a deck-spec JSON into a .pptx presentation.

Usage:
    python scripts/build_pptx.py <deckspec.json> [--out deliverables/<name>.pptx]
                                 [--charts deliverables/charts]

Template resolution: meta.template in templates/ -> first templates/*.pptx ->
built-in neutral 16:9 style. See specs/deck-spec.schema.md.
Run build_charts.py first so referenced chart PNGs exist.
"""
import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
RULE = RGBColor(0xE1, 0xE0, 0xD9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
BODY_TOP = Inches(1.55)

warnings = []


def resolve_template(meta, templates_dir):
    name = meta.get("template")
    if name:
        p = templates_dir / name
        if p.exists():
            return p
        warnings.append(f"template '{name}' not found in {templates_dir}; falling back")
    for p in sorted(templates_dir.glob("*.pptx")):
        return p
    return None


def new_presentation(meta, templates_dir):
    tpl = resolve_template(meta, templates_dir)
    if tpl:
        prs = Presentation(str(tpl))
        # Drop any existing slides so only spec content ships.
        xml_slides = prs.slides._sldIdLst  # noqa: SLF001 — python-pptx has no delete API
        for sld in list(xml_slides):
            xml_slides.remove(sld)
        print(f"Using template: {tpl.name}")
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        print("Using built-in neutral style")
    return prs


def blank_slide(prs):
    # Highest-index layout is typically Blank; fall back to the last one.
    layouts = prs.slide_layouts
    blank = None
    for lo in layouts:
        if lo.name.lower() == "blank":
            blank = lo
            break
    return prs.slides.add_slide(blank or layouts[len(layouts) - 1])


def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def set_run(p, text, size, color, bold=False, italic=False):
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return r


def add_headline(slide, text):
    _box, tf = textbox(slide, MARGIN, Inches(0.45), SLIDE_W - MARGIN * 2, Inches(0.95))
    set_run(tf.paragraphs[0], text, 24, INK, bold=True)
    # accent rule under headline
    line = slide.shapes.add_shape(1, MARGIN, Inches(1.38), Inches(1.6), Emu(28575))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_notes(slide, spec):
    notes = spec.get("notes", "")
    sources = spec.get("sources", [])
    if sources:
        notes = (notes + "\n" if notes else "") + "Sources: " + "; ".join(sources)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def check_bullets(spec, bullets):
    if len(bullets) > 6:
        warnings.append(f'slide "{spec.get("headline", spec["type"])}": {len(bullets)} bullets (>6)')
    for b in bullets:
        text = b["text"] if isinstance(b, dict) else b
        if len(re.findall(r"\S+", text)) > 12:
            warnings.append(f'bullet >12 words: "{text[:50]}..."')


def fill_bullets(tf, bullets, size=16):
    first = True
    for b in bullets:
        text, subs = (b.get("text", ""), b.get("sub", [])) if isinstance(b, dict) else (b, [])
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        set_run(p, f"• {text}", size, INK_2)
        p.space_after = Pt(10)
        for s in subs:
            sp = tf.add_paragraph()
            set_run(sp, f"– {s}", size - 2, MUTED)
            sp.level = 1
            sp.space_after = Pt(4)


def s_title(prs, spec, meta):
    slide = blank_slide(prs)
    _box, tf = textbox(slide, MARGIN, Inches(2.5), SLIDE_W - MARGIN * 2, Inches(1.4))
    set_run(tf.paragraphs[0], meta.get("title", "Untitled"), 38, INK, bold=True)
    sub = meta.get("subtitle")
    if sub:
        p = tf.add_paragraph()
        set_run(p, sub, 20, INK_2)
    _box2, tf2 = textbox(slide, MARGIN, Inches(6.4), SLIDE_W - MARGIN * 2, Inches(0.6))
    byline = " · ".join(x for x in [meta.get("author"), meta.get("date")] if x)
    if byline:
        set_run(tf2.paragraphs[0], byline, 12, MUTED)
    add_notes(slide, spec)


def s_agenda(prs, spec, _meta):
    slide = blank_slide(prs)
    add_headline(slide, spec.get("headline", "Agenda"))
    _box, tf = textbox(slide, MARGIN, BODY_TOP, SLIDE_W - MARGIN * 2, Inches(5.2))
    for i, item in enumerate(spec.get("items", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_run(p, f"{i + 1}   {item}", 18, INK_2)
        p.space_after = Pt(14)
    add_notes(slide, spec)


def s_content(prs, spec, _meta):
    slide = blank_slide(prs)
    add_headline(slide, spec["headline"])
    bullets = spec.get("bullets", [])
    check_bullets(spec, bullets)
    _box, tf = textbox(slide, MARGIN, BODY_TOP, SLIDE_W - MARGIN * 2, Inches(5.2))
    fill_bullets(tf, bullets)
    add_notes(slide, spec)


def s_two_column(prs, spec, _meta):
    slide = blank_slide(prs)
    add_headline(slide, spec["headline"])
    col_w = (SLIDE_W - MARGIN * 2 - Inches(0.4)) / 2
    for i, side in enumerate(("left", "right")):
        col = spec.get(side, {})
        x = MARGIN + i * (col_w + Inches(0.4))
        _box, tf = textbox(slide, x, BODY_TOP, col_w, Inches(5.2))
        set_run(tf.paragraphs[0], col.get("title", ""), 17, ACCENT, bold=True)
        tf.paragraphs[0].space_after = Pt(10)
        for b in col.get("bullets", []):
            p = tf.add_paragraph()
            text = b["text"] if isinstance(b, dict) else b
            set_run(p, f"• {text}", 15, INK_2)
            p.space_after = Pt(8)
    add_notes(slide, spec)


def s_chart(prs, spec, _meta, charts_dir):
    slide = blank_slide(prs)
    add_headline(slide, spec["headline"])
    png = charts_dir / f'{spec["chart_id"]}.png'
    if not png.exists():
        raise FileNotFoundError(f"chart PNG missing: {png} (run build_charts.py first)")
    pic_w = Inches(9.8)
    slide.shapes.add_picture(str(png), (SLIDE_W - pic_w) // 2, Inches(1.7), width=pic_w)
    takeaway = spec.get("takeaway")
    if takeaway:
        _box, tf = textbox(slide, MARGIN, Inches(6.85), SLIDE_W - MARGIN * 2, Inches(0.5))
        set_run(tf.paragraphs[0], takeaway, 13, INK_2, italic=True)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_notes(slide, spec)


def s_table(prs, spec, _meta):
    slide = blank_slide(prs)
    add_headline(slide, spec["headline"])
    cols, rows = spec["columns"], spec["rows"]
    shape = slide.shapes.add_table(len(rows) + 1, len(cols), MARGIN, BODY_TOP,
                                   SLIDE_W - MARGIN * 2,
                                   Inches(min(5.2, 0.45 * (len(rows) + 1))))
    table = shape.table
    for j, c in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = str(c)
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(v)
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    add_notes(slide, spec)


def s_swot(prs, spec, meta, charts_dir):
    # SWOT reuses the chart PNG for consistency with the standalone chart form.
    s_chart(prs, {**spec, "takeaway": spec.get("takeaway")}, meta, charts_dir)


def s_closing(prs, spec, _meta):
    slide = blank_slide(prs)
    _box, tf = textbox(slide, MARGIN, Inches(2.2), SLIDE_W - MARGIN * 2, Inches(1.2))
    set_run(tf.paragraphs[0], spec["headline"], 28, INK, bold=True)
    cta = spec.get("call_to_action")
    if cta:
        p = tf.add_paragraph()
        set_run(p, cta, 18, ACCENT, bold=True)
        p.space_before = Pt(12)
    steps = spec.get("next_steps", [])
    if steps:
        _b2, tf2 = textbox(slide, MARGIN, Inches(4.2), SLIDE_W - MARGIN * 2, Inches(2.6))
        set_run(tf2.paragraphs[0], "Next steps", 15, MUTED, bold=True)
        for s in steps:
            p = tf2.add_paragraph()
            set_run(p, f"• {s}", 15, INK_2)
            p.space_after = Pt(6)
    add_notes(slide, spec)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("spec")
    ap.add_argument("--out", default=None)
    ap.add_argument("--charts", default="deliverables/charts")
    ap.add_argument("--templates", default="templates")
    args = ap.parse_args()

    spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
    meta = spec.get("meta", {})
    charts_dir = Path(args.charts)
    chart_ids = {c["id"] for c in spec.get("charts", [])}

    # fail loudly on dangling chart refs before rendering anything
    for s in spec.get("slides", []):
        cid = s.get("chart_id")
        if cid and cid not in chart_ids:
            sys.exit(f'ERROR: slide "{s.get("headline", s["type"])}" references unknown chart_id "{cid}"')

    prs = new_presentation(meta, Path(args.templates))

    for s in spec.get("slides", []):
        t = s["type"]
        if t == "title":
            s_title(prs, s, meta)
        elif t == "agenda":
            s_agenda(prs, s, meta)
        elif t == "content":
            s_content(prs, s, meta)
        elif t == "two_column":
            s_two_column(prs, s, meta)
        elif t == "chart":
            s_chart(prs, s, meta, charts_dir)
        elif t == "table":
            s_table(prs, s, meta)
        elif t == "swot":
            s_swot(prs, s, meta, charts_dir)
        elif t == "closing":
            s_closing(prs, s, meta)
        else:
            sys.exit(f"ERROR: unknown slide type '{t}'")

    out = args.out
    if not out:
        safe = re.sub(r"[^\w\- ]", "", meta.get("title", "deck")).strip().replace(" ", "-")
        out = f"deliverables/{safe or 'deck'}.pptx"
    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Saved {out} ({len(prs.slides._sldIdLst)} slides)")  # noqa: SLF001

    for w in warnings:
        print(f"WARN: {w}", file=sys.stderr)


if __name__ == "__main__":
    main()
