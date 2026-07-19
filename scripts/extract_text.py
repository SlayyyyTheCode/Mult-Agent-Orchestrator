#!/usr/bin/env python3
"""Extract raw text from pptx / docx / pdf into location-tagged plain text.

Usage:
    python scripts/extract_text.py <file> [<file> ...]

Output (stdout): per file, markers the ingest agent turns into [src:] tags —
    === FILE: deck.pptx ===
    --- slide 1 ---
    ...text...
Unsupported or unreadable parts are marked [UNREADABLE: reason] so the agent
can flag them (L).
"""
import sys
from pathlib import Path

# Windows consoles default to cp1252 — force UTF-8 so em-dashes etc. survive.
sys.stdout.reconfigure(encoding="utf-8", errors="replace")


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, 1):
        print(f"--- slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(r.text for r in para.runs).strip()
                    if text:
                        print(text)
            elif shape.has_table:
                for row in shape.table.rows:
                    print(" | ".join(c.text.strip() for c in row.cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                print(f"[notes] {notes}")


def extract_docx(path):
    from docx import Document
    doc = Document(str(path))
    for i, para in enumerate(doc.paragraphs, 1):
        text = para.text.strip()
        if text:
            style = para.style.name if para.style else ""
            prefix = f"[{style}] " if style.startswith("Heading") else ""
            print(f"{prefix}{text}")
    for t_i, table in enumerate(doc.tables, 1):
        print(f"--- table {t_i} ---")
        for row in table.rows:
            print(" | ".join(c.text.strip() for c in row.cells))


def extract_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    for i, page in enumerate(reader.pages, 1):
        print(f"--- page {i} ---")
        text = (page.extract_text() or "").strip()
        print(text if text else "[UNREADABLE: no extractable text on this page]")


def extract_txt(path):
    print(Path(path).read_text(encoding="utf-8", errors="replace"))


HANDLERS = {
    ".pptx": extract_pptx,
    ".docx": extract_docx,
    ".pdf": extract_pdf,
    ".txt": extract_txt,
    ".md": extract_txt,
    ".vtt": extract_txt,
    ".srt": extract_txt,
}


def main():
    if len(sys.argv) < 2:
        sys.exit(__doc__)
    exit_code = 0
    for arg in sys.argv[1:]:
        path = Path(arg)
        print(f"=== FILE: {path.name} ===")
        if not path.exists():
            print("[UNREADABLE: file not found]")
            exit_code = 1
            continue
        handler = HANDLERS.get(path.suffix.lower())
        if not handler:
            print(f"[UNREADABLE: unsupported extension {path.suffix} — "
                  "images: read directly; audio: not supported in Stage 1]")
            continue
        try:
            handler(path)
        except Exception as e:  # noqa: BLE001 — report and continue with other files
            print(f"[UNREADABLE: {e}]")
            exit_code = 1
        print()
    sys.exit(exit_code)


if __name__ == "__main__":
    main()
